from django.shortcuts import render
from django.core.files.storage import FileSystemStorage
from django.contrib import messages
from django.http import HttpResponseRedirect
from django.core.mail import send_mail
from django.shortcuts import render, redirect
from django.db.models import Sum
import re
from django.utils import timezone
import os
import time
import openai
from pptx import Presentation
from django.http import HttpResponse
from io import BytesIO
from pptx.util import Pt
from pptx.dml.color import RGBColor
from django.http import HttpResponse, FileResponse
import base64
from django.contrib.auth.decorators import login_required
from django.http import HttpResponse
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph
from django.shortcuts import render
from docx import Document
from dotenv import load_dotenv
from .models import Project, ProjectBuy, JobApplication
from django.views.generic import DetailView


class ProjectDetailView(DetailView):
    model = Project
    template_name = 'project_detail.html'  # Specify the template name
    context_object_name = 'project'  # Specify the context object name


load_dotenv()


openai.api_key = os.getenv("OPENAI_API_KEY")


def home(request):
    return render(request, "index.html")


def about(request):
    return render(request, "about.html")


def contact(request):
    return render(request, "contact.html")


def service(request):
    return render(request, "service.html")


@login_required
def community(request):
    projects = Project.objects.all()
    if request.method == "GET":
        filter_criteria = request.GET.get('filter', None)
        if filter_criteria == 'website':
            projects = projects.filter(type='website')
        elif filter_criteria == 'android':
            projects = projects.filter(type='android')
        elif filter_criteria == 'machineLearning':
            projects = projects.filter(type='machineLearning')
        else:
            projects = Project.objects.all()

    return render(request, "community.html", {'projects': projects})


class ProjectDetail(DetailView):
    model = Project
    template_name = 'project_details.html'
    context_object_name = 'projects'

    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        project = Project.objects.all()
        context['project'] = project
        return context


def generate_pptx(request):
    return render(request, "generate-pptx.html")


@login_required
def generate_pptx_upload(request):
    if request.method == "POST":
        title = request.POST.get('title', '')
        description = request.POST.get('description', '')

        # Access the latest Text-Davinci model for content generation based on the provided title and description
        response = openai.Completion.create(
            engine="gpt-3.5-turbo-instruct",
            prompt=f"Title: {title}\nDescription: {description}\nWrite in detail, creating a comprehensive presentation with a minimum of 15 slides. Each slide should contain elaborate information and a distinct title.",
            max_tokens=3999  # Adjust max tokens for longer and more detailed content
        )
        generated_content = response['choices'][0]['text']

        # Create a presentation
        presentation = Presentation()

        # Set specific RGB color values for a refined color scheme
        colors = {
            'title': RGBColor(42, 87, 141),      # Title color
            'content': RGBColor(0, 0, 0),       # Content color (black)
            # Background color (light gray)
            'background': RGBColor(242, 242, 242)
        }

        # Process AI-generated content and create slides
        content_segments = generated_content.split("\n\n")
        num_slides = max(15, len(content_segments))

        for i in range(num_slides):
            # Use appropriate layout
            slide_layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(slide_layout)

            # Set slide background color
            background = slide.background
            background.fill.solid()
            background.fill.fore_color.rgb = colors['background']

            # Set original slide title
            current_title = f"Slide {i+1}: {title}" if i < len(
                content_segments) else f"Slide {i+1}: Original Slide Title"
            slide.shapes.title.text = current_title
            title_shape = slide.shapes.title
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.size = Pt(
                14)  # Larger font size for title
            title_shape.text_frame.paragraphs[0].font.color.rgb = colors['title']
            # Font style Arial
            title_shape.text_frame.paragraphs[0].font.name = 'Arial'

            # Set content text
            if i < len(content_segments):
                # Adjust placeholder index based on layout
                content_placeholder = slide.placeholders[1]
                content_placeholder.text = content_segments[i]
                content_shape = content_placeholder.text_frame
                content_shape.paragraphs[0].font.size = Pt(
                    14)  # Larger font size for content
                content_shape.paragraphs[0].font.color.rgb = colors['content']
                # Font style Calibri
                content_shape.paragraphs[0].font.name = 'Calibri'

        # Save the presentation in memory
        pptx_buffer = BytesIO()
        presentation.save(pptx_buffer)
        pptx_buffer.seek(0)

        encoded_pptx = base64.b64encode(pptx_buffer.getvalue()).decode('utf-8')

        context = {
            'pptx_file': encoded_pptx,
        }

        return render(request, 'generate-pptx.html', context)
    return render(request, 'generate-pptx.html')


@login_required
def generate_pdf_upload(request):
    if request.method == "POST":
        title = request.POST.get('title', '')
        description = request.POST.get('description', '')

        # Function to generate paragraphs with titles using OpenAI's Text-Davinci model
        def generate_paragraphs(title, length):
            response = openai.Completion.create(
                engine="gpt-3.5-turbo-instruct",
                prompt=f"Title: {title}\nDescription: {description}\nGenerate longer content with titles for the assignment PDF. Each section should span {length} pages.",
                max_tokens=3999  # Adjust max tokens for longer and more detailed content
            )
            return response['choices'][0]['text']

        # Generate paragraphs of specified lengths
        content = ""
        for i in range(5, 16):  # Generate paragraphs spanning 5 to 15 pages
            content += f"\n\nTitle: {title}\n\n"
            content += generate_paragraphs(title, i)

        # Create a PDF
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)

        # Set up styles for the document
        styles = getSampleStyleSheet()
        elements = []

        # Split the content into paragraphs and add to the PDF
        paragraphs = content.split("\n\n")
        for para in paragraphs:
            p = Paragraph(para, styles['Normal'])
            elements.append(p)

        doc.build(elements)
        buffer.seek(0)

        # Provide the generated PDF file as a download response without permanently saving it
        encoded_pdf_content = base64.b64encode(
            buffer.getvalue()).decode('utf-8')

        # Provide the generated PDF content as context to the template
        context = {
            'pdf_content': encoded_pdf_content,
        }

        return render(request, 'generate-pdf.html', context)

    return render(request, 'generate-pdf.html')


@login_required
def generate_doc_upload(request):
    if request.method == "POST":
        title = request.POST.get('title', '')
        description = request.POST.get('description', '')

        # Function to generate paragraphs using OpenAI's Text-Davinci model
        def generate_paragraphs(title, length):
            # Replace this with your actual OpenAI content generation logic
            # Use OpenAI API to generate content based on 'title' and 'length'
            # For example:
            response = openai.Completion.create(
                engine="gpt-3.5-turbo-instruct",
                prompt=f"Title: {title}\nDescription: {description}\nGenerate content with OpenAI.",
                max_tokens=100  # Adjust max_tokens as needed
            )
            return response['choices'][0]['text']

        # Generate content for the document
        content = ""
        for i in range(5, 16):  # Generate paragraphs spanning 5 to 15 pages
            content += f"\n\nTitle: {title}\n\n"
            content += generate_paragraphs(title, i)

        # Create a Word document
        buffer = BytesIO()
        doc = Document()

        # Split the content into paragraphs and add to the Word document
        paragraphs = content.split("\n\n")
        for para in paragraphs:
            doc.add_paragraph(para)

        # Save the document content to the buffer
        doc.save(buffer)
        buffer.seek(0)

        # Encode the document content to base64 for download
        encoded_doc_content = base64.b64encode(
            buffer.getvalue()).decode('utf-8')

        # Provide the generated Word content as context to the template
        context = {
            'doc_content': encoded_doc_content,
        }

        return render(request, 'generate-doc.html', context)

    return render(request, 'generate-doc.html')


@login_required
def fix_your_grammar(request):
    if request.method == "POST":
        user_input = request.POST.get('user_paragraph', '')

        # Instruction for grammar correction is included in the user's message
        user_message = (
            "Please review the following text for grammatical errors while ensuring "
            "the meaning remains unchanged. Correct any grammatical mistakes without "
            "altering the original context or style of this Paragraph.\n\n" + user_input
        )

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": user_message}],
            max_tokens=1900  # Adjust max tokens for longer and more detailed content
        )

        # Access the fixed paragraph from the API response
        fixed_text = response.choices[0].message['content']

        return render(request, 'grammar.html', {'fixed_text': fixed_text})

    return render(request, 'grammar.html')


@login_required
def paraphase(request):
    if request.method == "POST":
        user_input = request.POST.get('user_paragraph', '')

        # Instruction for grammar correction is included in the user's message
        user_message = (
            "Please rephrase each line of the provided text using simple, everyday language while keeping the original meaning and structure intact. Maintain the same number of lines in the output as in the input. The aim is to make small changes for uniqueness without altering the original meaning. Use common words for easy comprehension.\n\n" + user_input
        )

        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": user_message}],
            max_tokens=1900  # Adjust max tokens for longer and more detailed content
        )

        # Access the fixed paragraph from the API response
        fixed_text = response.choices[0].message['content']

        return render(request, 'paraphase.html', {'fixed_text': fixed_text})

    return render(request, 'paraphase.html')


@login_required
def visualize(request):
    if request.method == "POST":
        type = request.POST['type']
        company_description = request.POST.get('description', '')

        user_message = (
            f"Create {type} image. Based on {company_description}"
        )

        total_images_to_generate = 19
        generated_logos = []

        while total_images_to_generate > 0:
            # Limit the number of images generated in each API call to 10
            images_to_generate = min(total_images_to_generate, 10)

            response = openai.Image.create(
                prompt=user_message,
                n=images_to_generate,
            )

            generated_logos.extend(response['data'])
            total_images_to_generate -= images_to_generate

        return render(request, 'generate_logo.html', {'generated_logos': generated_logos})

    return render(request, 'generate_logo.html')





def upload_your_project(request):
    if request.method == "POST" or request.method == "FILE":
        name = request.POST['name']
        username = request.POST['username']
        email = request.POST['email']
        type = request.POST['type']
        price = request.POST['price']
        description = request.POST['description']
        photo = request.FILES['photo']
        project_database = Project(name=name, username=username, email=email,
                                   type=type, price=price, description=description, photo=photo)
        project_database.save()
        return HttpResponseRedirect(request.META.get('HTTP_REFERER'))
    return HttpResponseRedirect(request.META.get('HTTP_REFERER'))


def buy_project(request):
    if request.method == "POST":
        buyer_username = request.POST.get('buyer_username')
        project = request.POST.get('project')
        entry_exists = ProjectBuy.objects.filter(
            buyer_username=buyer_username, project=project).exists()

        if entry_exists:
            messages.error(
                request, "You Have Already Requested To Buy This.")
            return HttpResponseRedirect(request.META.get('HTTP_REFERER'))
        else:
            buyer_name = request.POST.get('buyer_name')
            buyer_email = request.POST.get('buyer_email')
            seller_name = request.POST.get('seller_name')
            seller_username = request.POST.get('seller_username')
            seller_email = request.POST.get('seller_email')
            type = request.POST.get('type')
            price = request.POST.get('price')
            description = request.POST.get('description')
            photo = request.POST.get('photo')

            project_buy_database = ProjectBuy(
                buyer_name=buyer_name,
                buyer_username=buyer_username,
                buyer_email=buyer_email,
                seller_name=seller_name,
                seller_username=seller_username,
                seller_email=seller_email,
                type=type,
                price=price,
                description=description,
                photo=photo,
                project=project
            )
            project_buy_database.save()
            messages.success(
                request, "Thank You For Your Request To Buy This. We Will Contact With The Seller & Let You Know.")
            return HttpResponseRedirect(request.META.get('HTTP_REFERER'))
    return HttpResponseRedirect(request.META.get('HTTP_REFERER'))


def job_application(request):
    if request.method == 'POST':
        category = request.POST.get('category')
        name = request.POST.get('name')
        email = request.POST.get('email')
        phone = request.POST.get('phone')
        country = request.POST.get('country')
        cv = request.FILES.get('cv')
        photo = request.FILES.get('photo')
        education = request.POST.get('education')
        experience = request.POST.get('experience')

        # Create a JobApplication instance
        job_application = JobApplication(
            category=category,
            name=name,
            email=email,
            phone=phone,
            country=country,
            cv=cv,
            photo=photo,
            education=education,
            experience=experience
        )
        job_application.save()
        messages.success(
                request, "Thank you for your interest in Think AI careers. We'll review your CV and reach out soon if there's a match. Regards, Think AI Team.")
        # Email details
        sender_email = 'thinkai@thinkai.fun'
        receiver_email = email
        subject = 'Thank You For Applying In Our Career'
        message = 'Dear ' + name + \
            ',\n\nWe have received your application. We will review your CV and reach out soon if there is a match.\n\nRegards\nThink AI Team'

        # Send email
        send_mail(
            subject,
            message,
            sender_email,
            [receiver_email],
            fail_silently=False,
        )
        return HttpResponseRedirect(request.META.get('HTTP_REFERER'))
    return render(request, 'career.html')

